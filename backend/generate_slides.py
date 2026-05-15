"""Generate presentation slides as .pptx for the drone navigation project."""

import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ASSETS = pathlib.Path(__file__).parent.parent / "docs" / "assets"
OUT = pathlib.Path(__file__).parent.parent / "docs" / "drone-nav-presentation.pptx"

# ── Color palette ────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x1E, 0x1E, 0x2E)   # dark navy
ACCENT    = RGBColor(0x89, 0xDC, 0xEB)   # cyan
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW    = RGBColor(0xFF, 0xD7, 0x00)
GREEN     = RGBColor(0x4C, 0xAF, 0x50)
RED       = RGBColor(0xE5, 0x39, 0x35)
GRAY_LIGHT= RGBColor(0xCC, 0xCC, 0xCC)
DEMO_BG   = RGBColor(0x0D, 0x1B, 0x2A)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

def slide():
    s = prs.slides.add_slide(BLANK)
    # Dark background
    bg = s.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return s

def txb(slide, text, l, t, w, h, size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf

def heading(s, text, t=0.3):
    txb(s, text, 0.4, t, 12.5, 0.7, size=32, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

def divider(s, t):
    from pptx.util import Pt as Pt2
    ln = s.shapes.add_shape(1, Inches(0.4), Inches(t), Inches(12.5), Inches(0.04))
    ln.fill.solid()
    ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background()

def bullet_block(s, items, l, t, w, h, size=18, color=WHITE, title=None, title_color=None):
    tf = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = True
    frame = tf.text_frame
    frame.word_wrap = True
    first = True
    if title:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color or ACCENT
        first = False
    for item in items:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        first = False

def table(s, headers, rows, l, t, w, h):
    from pptx.util import Pt
    cols = len(headers)
    tbl = s.shapes.add_table(len(rows)+1, cols,
                              Inches(l), Inches(t), Inches(w), Inches(h)).table
    tbl.columns[0].width = Inches(w * 0.40)
    for i in range(1, cols):
        tbl.columns[i].width = Inches(w * 0.20)

    def cell_style(cell, text, bold=False, bg=None, fg=WHITE, sz=14, align=PP_ALIGN.CENTER):
        cell.text = text
        cell.text_frame.paragraphs[0].alignment = align
        run = cell.text_frame.paragraphs[0].runs[0] if cell.text_frame.paragraphs[0].runs else cell.text_frame.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = fg
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    for c, h in enumerate(headers):
        cell_style(tbl.cell(0, c), h, bold=True,
                   bg=RGBColor(0x2A,0x2A,0x4A), fg=ACCENT, sz=13)

    row_colors = [RGBColor(0x25,0x25,0x3A), RGBColor(0x1E,0x1E,0x2E)]
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            fg = YELLOW if val in ("⚡ Bounded", "⚠️ Not always") else (GREEN if "✅" in val else WHITE)
            cell_style(tbl.cell(r+1, c), val,
                       bg=row_colors[r % 2], fg=fg, sz=13, align=align)

def add_image(s, path, l, t, w, h):
    p = pathlib.Path(path)
    if p.exists():
        s.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        # placeholder box
        ph = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        ph.fill.solid(); ph.fill.fore_color.rgb = RGBColor(0x33,0x33,0x44)
        ph.line.color.rgb = ACCENT
        tf = ph.text_frame
        tf.text = f"[image: {p.name}]"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].runs[0].font.color.rgb = GRAY_LIGHT
        tf.paragraphs[0].runs[0].font.size = Pt(12)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
s1 = slide()
txb(s1, "AI-Based Autonomous Drone Navigation",
    0.5, 1.8, 12.3, 1.2, size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
txb(s1, "in a Simulated 3D Urban Environment",
    0.5, 2.85, 12.3, 0.8, size=28, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
divider(s1, 3.8)
txb(s1, "COMP 569 – Artificial Intelligence  |  Prof. Reza Abdolee",
    0.5, 4.0, 12.3, 0.5, size=18, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)
txb(s1, "Swathi Krishna Cheripally  |  May 2026",
    0.5, 4.6, 12.3, 0.5, size=18, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Motivation & Problem
# ─────────────────────────────────────────────────────────────────────────────
s2 = slide()
heading(s2, "Motivation & Problem")
divider(s2, 1.1)

txb(s2, "Why does this matter?", 0.5, 1.3, 6.0, 0.4, size=18, bold=True, color=ACCENT)
bullet_block(s2, [
    "🚁  Package delivery, inspection, emergency response",
    "🏙️  Must avoid buildings AND restricted no-fly zones",
    "⚡  Climbing costs MORE energy than horizontal movement",
    "🎯  Goal: find the lowest-energy path from A → B",
], 0.5, 1.7, 6.2, 2.8, size=17)

# Right side: cost model box
box = s2.shapes.add_shape(1, Inches(7.2), Inches(1.3), Inches(5.7), Inches(4.8))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x25,0x25,0x3A)
box.line.color.rgb = ACCENT

txb(s2, "Movement Cost Model", 7.4, 1.4, 5.3, 0.5, size=17, bold=True, color=ACCENT)
cost_rows = [
    ("← → ↑↓ (horizontal)", "cost  1"),
    ("↓  downward",          "cost  1"),
    ("↑  upward",            "cost  2  ⚠️"),
]
for i, (direction, cost) in enumerate(cost_rows):
    txb(s2, direction, 7.4, 2.0 + i*0.6, 3.5, 0.5, size=16, color=WHITE)
    txb(s2, cost,      11.0, 2.0 + i*0.6, 1.7, 0.5, size=16,
        color=YELLOW if "2" in cost else WHITE)

txb(s2, "State:  s = (x, y, z)  —  bounded 3D grid", 7.4, 3.8, 5.3, 0.5, size=15, color=GRAY_LIGHT)
txb(s2, "Obstacles + No-Fly Zones = impassable cells", 7.4, 4.3, 5.3, 0.5, size=15, color=GRAY_LIGHT)
txb(s2, "Formally: state-space search  <S, A, T, C, G>", 7.4, 4.8, 5.3, 0.5, size=15, color=GRAY_LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — System Architecture
# ─────────────────────────────────────────────────────────────────────────────
s3 = slide()
heading(s3, "System Architecture")
divider(s3, 1.1)

# Frontend box
fb = s3.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(5.2), Inches(4.5))
fb.fill.solid(); fb.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x1A)
fb.line.color.rgb = GREEN

txb(s3, "Frontend", 0.7, 1.6, 4.8, 0.5, size=20, bold=True, color=GREEN)
txb(s3, "React + TypeScript + Vite", 0.7, 2.1, 4.8, 0.4, size=14, color=GRAY_LIGHT)
bullet_block(s3, [
    "🗺  Three.js 3D visualization",
    "🎛  Algorithm selector & coordinate editor",
    "🏗  Obstacle & no-fly zone controls",
    "✈️  Animated drone along computed path",
    "📊  Live metrics display",
], 0.7, 2.5, 4.8, 3.0, size=15)

# Arrow
txb(s3, "POST /pathfind\nJSON over HTTP", 5.9, 3.1, 1.5, 1.0, size=13, color=ACCENT, align=PP_ALIGN.CENTER)
arr = s3.shapes.add_shape(1, Inches(5.8), Inches(3.6), Inches(1.7), Inches(0.08))
arr.fill.solid(); arr.fill.fore_color.rgb = ACCENT
arr.line.fill.background()

# Backend box
bb = s3.shapes.add_shape(1, Inches(7.6), Inches(1.5), Inches(5.2), Inches(4.5))
bb.fill.solid(); bb.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
bb.line.color.rgb = ACCENT

txb(s3, "Backend", 7.8, 1.6, 4.8, 0.5, size=20, bold=True, color=ACCENT)
txb(s3, "Python 3 + FastAPI + Pydantic", 7.8, 2.1, 4.8, 0.4, size=14, color=GRAY_LIGHT)
bullet_block(s3, [
    "🔌  POST /pathfind  — single endpoint",
    "🏙  Environment model (3D grid + costs)",
    "🧠  6 search algorithms",
    "📐  Pydantic request/response validation",
    "🔓  CORS enabled for frontend access",
], 7.8, 2.5, 4.8, 3.0, size=15)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — The 6 Algorithms
# ─────────────────────────────────────────────────────────────────────────────
s4 = slide()
heading(s4, "The 6 Search Algorithms")
divider(s4, 1.1)

headers = ["Algorithm", "Key Idea", "Optimal?"]
rows = [
    ("UCS", "Expands by cost only — no goal awareness  f(n)=g(n)", "✅ Yes"),
    ("A* Manhattan", "Cost + city-block distance to goal  f=g+h", "✅ Yes"),
    ("A* Euclidean", "Cost + straight-line distance to goal", "✅ Yes"),
    ("A* Building-Aware", "Accounts for obstacle height early (custom)", "⚠️ Not always"),
    ("IDA*", "DFS with increasing f-thresholds — O(depth) memory", "✅ Yes"),
    ("Weighted A* (w=1.5)", "Inflated heuristic: f = g + 1.5×h — rushes to goal", "⚡ Bounded"),
]
table(s4, headers, rows, 0.4, 1.3, 12.5, 5.5)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — LIVE DEMO
# ─────────────────────────────────────────────────────────────────────────────
s5 = slide()
# Full-screen demo background
bg2 = s5.background.fill
bg2.solid()
bg2.fore_color.rgb = DEMO_BG

txb(s5, "LIVE DEMO", 0.5, 1.8, 12.3, 1.2, size=56, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
txb(s5, "3D Drone Navigation — Interactive Frontend",
    0.5, 3.2, 12.3, 0.6, size=22, color=WHITE, align=PP_ALIGN.CENTER)
divider(s5, 4.1)

txb(s5, "http://localhost:5173", 0.5, 4.3, 12.3, 0.5,
    size=18, color=ACCENT, align=PP_ALIGN.CENTER)

bullet_block(s5, [
    "1.  Generate Random City  →  gray building columns appear",
    "2.  Generate No-Fly Zones  →  red restricted blocks appear",
    "3.  Select A* Manhattan  →  Find Path  →  read metrics",
    "4.  Switch to Weighted A*  →  same cost, far fewer nodes",
    "5.  Switch to IDA*  →  same path, ~11,000 nodes, ~118ms",
], 1.5, 4.9, 10.3, 2.3, size=16, color=GRAY_LIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Results
# ─────────────────────────────────────────────────────────────────────────────
s6 = slide()
heading(s6, "Benchmark Results — 7 Scenarios × 5 Algorithms")
divider(s6, 1.1)

# Chart image left
add_image(s6, ASSETS / "chart_nodes_expanded.png", 0.4, 1.3, 7.2, 4.5)

# Key numbers right
box2 = s6.shapes.add_shape(1, Inches(7.9), Inches(1.3), Inches(5.0), Inches(4.5))
box2.fill.solid(); box2.fill.fore_color.rgb = RGBColor(0x25,0x25,0x3A)
box2.line.color.rgb = YELLOW

txb(s6, "Key Numbers", 8.1, 1.4, 4.6, 0.5, size=18, bold=True, color=YELLOW)

stats = [
    ("Weighted A* vs UCS", "95–98% fewer nodes"),
    ("Weighted A* vs IDA*", "320× fewer nodes"),
    ("IDA* nodes", "11,227"),
    ("Weighted A* nodes", "35 (same env)"),
    ("All 5 algorithms", "same optimal cost"),
]
for i, (label, val) in enumerate(stats):
    txb(s6, label, 8.1, 2.0 + i*0.7, 2.8, 0.5, size=14, color=GRAY_LIGHT)
    txb(s6, val,   11.0, 2.0 + i*0.7, 1.8, 0.5, size=14, bold=True,
        color=YELLOW if "×" in val or "%" in val else (GREEN if "same" in val else WHITE))

txb(s6, "✅  Correctness confirmed — all algorithms agree on cost",
    7.9, 5.45, 5.0, 0.4, size=13, color=GREEN)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Screenshots comparison
# ─────────────────────────────────────────────────────────────────────────────
s7 = slide()
heading(s7, "IDA* vs Weighted A* — Same Environment")
divider(s7, 1.1)

img_dir = ASSETS / "drone-final-images"
add_image(s7, img_dir / "drone-final-6.png", 0.4, 1.3, 5.9, 4.5)
add_image(s7, img_dir / "drone-final-7.png", 6.9, 1.3, 5.9, 4.5)

txb(s7, "IDA*", 0.4, 5.85, 5.9, 0.4, size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)
txb(s7, "11,227 nodes  |  118 ms", 0.4, 6.25, 5.9, 0.4, size=14, color=WHITE, align=PP_ALIGN.CENTER)

txb(s7, "Weighted A* (w=1.5)", 6.9, 5.85, 5.9, 0.4, size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txb(s7, "35 nodes  |  1.4 ms  →  320× fewer", 6.9, 6.25, 5.9, 0.4, size=14, color=WHITE, align=PP_ALIGN.CENTER)

txb(s7, "Same cost: 15  |  Same path  |  Same environment",
    0.4, 6.9, 12.5, 0.4, size=15, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — Conclusion
# ─────────────────────────────────────────────────────────────────────────────
s8 = slide()
heading(s8, "Conclusion")
divider(s8, 1.1)

findings = [
    ("✅", "All 5 algorithms found identical optimal costs across 7 scenarios",
     "Correctness validated"),
    ("📉", "Heuristic guidance reduces node expansions by 30–98% vs UCS",
     "Informed search wins"),
    ("⚡", "Weighted A* (w=1.5): 95–98% fewer nodes, sub-ms runtime",
     "Practical winner"),
    ("🧠", "IDA* expands 320× more nodes than Weighted A* in 3D grids",
     "Impractical at scale"),
    ("⚠️", "Building-aware heuristic found inadmissible — needs redesign",
     "Heuristics need proof"),
]

for i, (icon, text, label) in enumerate(findings):
    y = 1.5 + i * 0.95
    box_f = s8.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Inches(0.8))
    box_f.fill.solid()
    box_f.fill.fore_color.rgb = RGBColor(0x25,0x25,0x3A)
    box_f.line.color.rgb = RGBColor(0x40,0x40,0x60)
    txb(s8, icon, 0.6, y+0.1, 0.6, 0.6, size=22, align=PP_ALIGN.CENTER)
    txb(s8, text, 1.3, y+0.12, 8.8, 0.6, size=15, color=WHITE)
    txb(s8, label, 10.2, y+0.12, 2.4, 0.6, size=13, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

txb(s8, "Thank you  —  Questions welcome",
    0.4, 6.9, 12.5, 0.4, size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"Slides saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
